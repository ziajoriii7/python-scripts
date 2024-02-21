import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
import os
import re
from tqdm import tqdm

parser = argparse.ArgumentParser(
    description="Create a PowerPoint presentation with images"
)
parser.add_argument(
    "-maximize",
    "-max",
    action="store_true",
    help="Maximize images size to full the slide, ideal for screenshots of other pptxs, can distort images otherwise.",
)
parser.add_argument(
    "-border", "border", type=int, help="Add a border around the image, use pixels."
)
args = parser.parse_args()


def atoi(text):
    return int(text) if text.isdigit() else text


def natural_keys(text):
    return [atoi(c) for c in re.split(r"(\d+)", text)]


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


output_file = input("name of the file? (default is presentation.pptx) ")
if not output_file:
    output_file = "presentation.pptx"

image_folder = input("Insert which directory has the pics: ")

image_folder = image_folder.strip('""')


image_files = sorted(
    [
        f
        for f in os.listdir(image_folder)
        if f.endswith((".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"))
    ],
    key=natural_keys,
)

for image_file in tqdm(image_files, desc="Processing", unit=" file"):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    image_path = os.path.join(image_folder, image_file)

    if args.maximize:
        if args.border is not None:
            border = Pt(args.border)
            img = slide.shapes.add_picture(
                image_path,
                border,
                border,
                width=prs.slide_width - 2 * border,
                height=prs.slide_height - 2 * border,
            )
        else:
            img = slide.shapes.add_picture(
                image_path, 0, 0, width=prs.slide_width, height=prs.slide_height
            )

    else:
        if args.border is not None:
            border = Pt(args.border)
            effective_height = prs.slide_height - 2 * border
            img = slide.shapes.add_picture(
                image_path, border, border, height=effective_height
            )
            ratio = effective_height / img.height
        else:
            img = slide.shapes.add_picture(image_path, 0, 0, height=prs.slide_height)
            ratio = prs.slide_height / img.height

        img.width = int(img.width * ratio)
        img.left = int((prs.slide_width - img.width) / 2)
        img.top = int((prs.slide_height - img.height) / 2)


prs.save(output_file)
print(f"Presentation saved to {output_file}")
